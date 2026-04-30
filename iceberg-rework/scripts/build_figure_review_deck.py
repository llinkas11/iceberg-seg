"""
Build the paper figure-review PowerPoint deck.

The deck is a review artifact: each slide contains only one rendered artifact
(a figure PNG or a table-preview PNG). Speaker notes carry the absolute
pathname first, then a concise caption. The companion CSV is the checklist.
"""

from __future__ import annotations

import argparse
import csv
import re
import tempfile
from collections import defaultdict
from pathlib import Path
from zipfile import ZIP_DEFLATED, ZipFile
from xml.sax.saxutils import escape

import matplotlib

matplotlib.use("Agg")
import matplotlib.pyplot as plt
from PIL import Image
from pptx import Presentation
from pptx.util import Inches


ROOT = Path(__file__).resolve().parents[2]
FIG_ARCHIVE = ROOT / "paper-writing/figures/fig-archive"
RESULTS_MD = ROOT / "paper-writing/results.md"
OVERLEAF_MAIN = ROOT / "paper-writing/overleaf/2026-04-17_template-plan-md/main.tex"

FIGURE_ORDER = [
    "fig01_annotation_difficulty",
    "fig02_dataset_workflow",
    "fig03_method_schematic",
    "fig04_evaluation_schematic",
    "fig05_progression",
    "mae_rootlen_vs_sza",
    "area_scatter_by_method",
    "bias_delta_by_area",
    "re_by_area_bin",
    "outline_examples",
]

FIGURE_CAPTIONS = {
    "fig01_annotation_difficulty": (
        "Annotation-difficulty examples for Sentinel-2 iceberg chips. The "
        "panels contrast the source image, preliminary annotation, cleaned "
        "binary target, and the main annotation decision introduced by the "
        "v4_clean workflow."
    ),
    "fig02_dataset_workflow": (
        "Dataset construction workflow. Fisser-derived and Roboflow-derived "
        "annotations are harmonized through shadow merge, 40 m component "
        "filtering, and annotation-aware IC masking to produce the binary "
        "v4_clean train, validation, and test splits."
    ),
    "fig03_method_schematic": (
        "Base segmentation workflow schematic. Two methods operate directly "
        "on B08 reflectance, while learned methods share a UNet++ probability "
        "map before direct segmentation, fixed-threshold binarization, Otsu "
        "binarization, or CRF refinement."
    ),
    "fig04_evaluation_schematic": (
        "Per-iceberg evaluation workflow. Ground-truth and predicted "
        "components are matched by Hungarian assignment on 1 - IoU, unmatched "
        "components are counted as false negatives or false positives, and "
        "matched pairs feed MAE, IoU, and MSE summaries."
    ),
    "fig05_progression": (
        "Experimental progression visual. Phase A walks dataset cleaning and "
        "balancing choices with one reader-facing motivation per step; Phase "
        "B freezes the selected dataset and compares base segmentation "
        "workflows, with top-hat variants treated as a separate downstream "
        "post-processing result."
    ),
    "mae_rootlen_vs_sza": (
        "Per-pair root-length MAE by SZA bin for the method comparison. "
        "Lower values indicate smaller matched-iceberg size error, making "
        "this the most direct Fisser-comparable accuracy view."
    ),
    "area_scatter_by_method": (
        "Predicted versus reference iceberg area by method. The scatter view "
        "shows method-specific area bias and spread across matched iceberg "
        "pairs."
    ),
    "bias_delta_by_area": (
        "Area-dependent bias summary by method. The figure shows whether "
        "prediction error changes across iceberg size, especially for small "
        "versus large objects."
    ),
    "re_by_area_bin": (
        "Relative error by iceberg area bin. Values compare predicted and "
        "reference area within size classes, exposing whether method effects "
        "concentrate in particular iceberg sizes."
    ),
    "outline_examples": (
        "Representative outline examples comparing predicted and reference "
        "iceberg boundaries. These examples provide visual context for the "
        "numeric MAE and IoU metrics."
    ),
}

TABLE_SPECS = [
    {
        "slug": "table_phase_a_leaderboard",
        "title": "Phase A leaderboard",
        "source": RESULTS_MD,
        "caption": (
            "Phase A lt65 dataset-progression leaderboard from results.md. "
            "A0 is the best validation-IoU configuration, while later "
            "cleaning and balancing variants underperform on this split."
        ),
        "headers": [
            "ID",
            "Manifest",
            "val IoU",
            "test IoU",
            "UNet match rate",
            "UNet RL MAE (m)",
        ],
        "rows": [
            ["A0", "v4_raw_lt65\nFisser preprocessing, no nulls", "0.613", "0.577", "0.512", "9.82"],
            ["A1", "v4_raw_lt65_plus_nulls\nFisser preprocessing + nulls", "0.503", "0.477", "0.315", "15.21"],
            ["A3", "v4_clean_lt65_plus_nulls", "0.269", "0.336", "0.182", "15.69"],
            ["A2", "v4_clean_lt65\nour preprocessing, no nulls", "0.261", "0.344", "0.245", "15.26"],
            ["A7-A9", "v4_clean + nulls + aug\n+ size oversample", "0.243", "0.320", "0.163", "14.78"],
            ["A5-A6", "v4_clean + nulls + aug\n+ class balance", "0.237", "0.312", "0.158", "15.23"],
            ["A4", "v4_clean + nulls + aug", "0.225", "0.274", "0.122", "14.93"],
        ],
    },
    {
        "slug": "table_phase_a_preprocessing_nulls",
        "title": "Phase A preprocessing x null-chip injection",
        "source": RESULTS_MD,
        "caption": (
            "Phase A 2x2 from results.md. Preprocessing dominates the lt65 "
            "validation-IoU change, while null-chip injection has a smaller "
            "effect."
        ),
        "headers": ["Preprocessing", "No nulls", "+ nulls 1:1"],
        "rows": [
            ["Fisser preprocessing", "0.613", "0.503"],
            ["Our preprocessing", "0.261", "0.269"],
        ],
    },
    {
        "slug": "table_probability_calibration",
        "title": "Probability calibration audit",
        "source": RESULTS_MD,
        "caption": (
            "Probability calibration audit from results.md. A2 shifts most "
            "pixels into the 0.20-0.35 probability range, explaining the "
            "failure of fixed-threshold post-processing on that run."
        ),
        "headers": ["Run", "median P(iceberg)", "frac. pixels 0.20-0.35", "frac. pixels >= 0.5"],
        "rows": [
            ["baseline_v1", "0.001", "0.4%", "4.0%"],
            ["A0 (raw lt65)", "0.013", "2.2%", "6.6%"],
            ["A2 (clean lt65)", "0.278", "59.0%", "5.2%"],
        ],
    },
    {
        "slug": "table_base_mae_rootlen",
        "title": "Base methods: root-length MAE (m)",
        "source": RESULTS_MD,
        "caption": (
            "Base six-method root-length MAE table from results.md. UNet_CRF "
            "has the lowest MAE in three of four SZA bins; UNet_OT wins the "
            "lt65 bin."
        ),
        "headers": ["Method", "< 65", "65-70", "70-75", "> 75"],
        "rows": [
            ["TR", "17.81", "7.91", "6.46", "20.07"],
            ["OT", "22.73", "13.77", "14.51", "15.91"],
            ["UNet", "10.48", "11.48", "13.86", "15.57"],
            ["UNet_TR", "14.24", "15.54", "18.56", "19.62"],
            ["UNet_OT", "7.98", "11.96", "13.91", "15.27"],
            ["UNet_CRF", "10.12", "7.37", "9.04", "12.59"],
        ],
    },
    {
        "slug": "table_base_iou",
        "title": "Base methods: per-pair IoU",
        "source": RESULTS_MD,
        "caption": (
            "Base six-method matched-pair IoU table from results.md. UNet_CRF "
            "again leads in three of four SZA bins, with UNet_OT leading lt65."
        ),
        "headers": ["Method", "< 65", "65-70", "70-75", "> 75"],
        "rows": [
            ["TR", "0.482", "0.670", "0.686", "0.594"],
            ["OT", "0.470", "0.622", "0.621", "0.620"],
            ["UNet", "0.701", "0.672", "0.643", "0.646"],
            ["UNet_TR", "0.665", "0.639", "0.603", "0.611"],
            ["UNet_OT", "0.730", "0.673", "0.643", "0.642"],
            ["UNet_CRF", "0.653", "0.691", "0.666", "0.658"],
        ],
    },
    {
        "slug": "table_base_detection",
        "title": "Base methods: detection statistics",
        "source": RESULTS_MD,
        "caption": (
            "Base six-method detection table from results.md. Match rate and "
            "precision disclose selection bias before comparing MAE across "
            "methods."
        ),
        "headers": ["Method", "n_gt", "n_pred", "n_matched", "match rate", "precision"],
        "rows": [
            ["TR", "18,990", "16,818", "2,759", "0.145", "0.164"],
            ["OT", "18,990", "35,564", "5,547", "0.292", "0.156"],
            ["UNet", "18,990", "21,422", "9,916", "0.522", "0.463"],
            ["UNet_TR", "18,990", "23,009", "8,369", "0.441", "0.364"],
            ["UNet_OT", "18,990", "15,468", "5,929", "0.312", "0.483"],
            ["UNet_CRF", "18,990", "20,981", "8,891", "0.468", "0.424"],
        ],
    },
    {
        "slug": "table_tophat_mae_rootlen",
        "title": "Top-hat variants: root-length MAE (m)",
        "source": OVERLEAF_MAIN,
        "caption": (
            "Top-hat post-processing root-length MAE table from Overleaf "
            "Table 9. These +TH variants are the separate Segmentation -> "
            "top-hat -> segmentation++ sensitivity path."
        ),
        "headers": ["Method", "< 65", "65-70", "70-75", "> 75"],
        "rows": [
            ["TR+TH", "17.1", "8.0", "11.0", "20.4"],
            ["OT+TH", "17.4", "12.6", "13.8", "19.5"],
            ["UNet+TH", "10.3", "12.0", "17.5", "23.2"],
            ["UNet_TR+TH", "13.9", "15.8", "20.4", "24.4"],
            ["UNet_OT+TH", "13.2", "12.5", "16.4", "22.0"],
            ["UNet_CRF+TH", "10.0", "8.5", "12.9", "20.7"],
        ],
    },
    {
        "slug": "table_tophat_iou",
        "title": "Top-hat variants: per-pair IoU",
        "source": OVERLEAF_MAIN,
        "caption": (
            "Top-hat post-processing IoU table from Overleaf Table 10. It "
            "keeps the +TH sensitivity branch separate from the six base "
            "methods."
        ),
        "headers": ["Method", "< 65", "65-70", "70-75", "> 75"],
        "rows": [
            ["TR+TH", "0.51", "0.66", "0.65", "0.60"],
            ["OT+TH", "0.52", "0.61", "0.62", "0.58"],
            ["UNet+TH", "0.72", "0.66", "0.64", "0.60"],
            ["UNet_TR+TH", "0.68", "0.64", "0.61", "0.59"],
            ["UNet_OT+TH", "0.62", "0.66", "0.64", "0.60"],
            ["UNet_CRF+TH", "0.67", "0.68", "0.66", "0.61"],
        ],
    },
    {
        "slug": "table_tophat_detection",
        "title": "Top-hat variants: detection statistics",
        "source": OVERLEAF_MAIN,
        "caption": (
            "Top-hat post-processing detection table from Overleaf Table 11. "
            "The +TH branch trades precision for additional recovered "
            "candidate matches."
        ),
        "headers": ["Method", "n_ref", "n_pred", "n_matched", "match rate (%)", "precision (%)"],
        "rows": [
            ["TR+TH", "7201", "28342", "2851", "39.6", "10.1"],
            ["OT+TH", "7201", "42859", "2883", "40.0", "6.7"],
            ["UNet+TH", "7201", "14174", "3147", "43.7", "22.2"],
            ["UNet_TR+TH", "7201", "14080", "2874", "39.9", "20.4"],
            ["UNet_OT+TH", "7201", "17630", "2963", "41.2", "16.8"],
            ["UNet_CRF+TH", "7201", "16065", "3330", "46.2", "20.7"],
        ],
    },
]

NOTES_SLIDE_REL = "http://schemas.openxmlformats.org/officeDocument/2006/relationships/notesSlide"
NOTES_MASTER_REL = "http://schemas.openxmlformats.org/officeDocument/2006/relationships/notesMaster"
SLIDE_REL = "http://schemas.openxmlformats.org/officeDocument/2006/relationships/slide"
REL_NS = "http://schemas.openxmlformats.org/package/2006/relationships"
NOTES_MASTER_CONTENT_TYPE = "application/vnd.openxmlformats-officedocument.presentationml.notesMaster+xml"
NOTES_SLIDE_CONTENT_TYPE = "application/vnd.openxmlformats-officedocument.presentationml.notesSlide+xml"
THEME_CONTENT_TYPE = "application/vnd.openxmlformats-officedocument.theme+xml"


def latest_pngs_by_slug(fig_dir):
    by_slug = defaultdict(list)
    for path in fig_dir.glob("*.png"):
        if "__" not in path.stem:
            continue
        timestamp, slug = path.stem.split("__", 1)
        by_slug[slug].append((timestamp, path))
    return {slug: sorted(values)[-1][1] for slug, values in by_slug.items()}


def render_table(spec, out_dir):
    out_dir.mkdir(parents=True, exist_ok=True)
    headers = spec["headers"]
    rows = spec["rows"]
    ncols = len(headers)
    nrows = len(rows)
    fig_w = max(8.0, 1.65 * ncols)
    fig_h = 1.8 + 0.52 * (nrows + 1)
    fig, ax = plt.subplots(figsize=(fig_w, fig_h))
    fig.patch.set_facecolor("white")
    ax.axis("off")
    ax.text(
        0.5,
        0.98,
        spec["title"],
        ha="center",
        va="top",
        fontsize=15,
        weight="bold",
        transform=ax.transAxes,
    )
    table = ax.table(
        cellText=rows,
        colLabels=headers,
        cellLoc="center",
        loc="center",
        bbox=[0.02, 0.04, 0.96, 0.82],
    )
    table.auto_set_font_size(False)
    table.set_fontsize(9)
    table.scale(1, 1.25)
    for (r, _c), cell in table.get_celld().items():
        cell.set_edgecolor("#333333")
        cell.set_linewidth(0.55)
        if r == 0:
            cell.set_facecolor("#E8EEF7")
            cell.get_text().set_weight("bold")
        elif r % 2 == 0:
            cell.set_facecolor("#F7F7F7")
    out_path = out_dir / f"{spec['slug']}.png"
    fig.savefig(out_path, dpi=220, bbox_inches="tight", facecolor="white")
    plt.close(fig)
    return out_path


def collect_artifacts(out_dir):
    latest = latest_pngs_by_slug(FIG_ARCHIVE)
    rows = []
    for slug in FIGURE_ORDER:
        if slug not in latest:
            raise FileNotFoundError(f"Missing figure slug in {FIG_ARCHIVE}: {slug}")
        png_path = latest[slug]
        svg_path = png_path.with_suffix(".svg")
        rows.append(
            {
                "slug": slug,
                "kind": "figure",
                "path": png_path,
                "svg_path": svg_path if svg_path.exists() else None,
                "caption": FIGURE_CAPTIONS[slug],
            }
        )

    table_dir = out_dir / "table_previews"
    for spec in TABLE_SPECS:
        rows.append(
            {
                "slug": spec["slug"],
                "kind": "table",
                "path": render_table(spec, table_dir),
                "svg_path": None,
                "caption": f"{spec['caption']} Source: {spec['source'].relative_to(ROOT)}.",
            }
        )
    return rows


def add_full_slide_picture(slide, artifact_path, prs):
    with Image.open(artifact_path) as im:
        px_w, px_h = im.size
    max_w = prs.slide_width - Inches(0.45)
    max_h = prs.slide_height - Inches(0.45)
    scale = min(max_w / px_w, max_h / px_h)
    pic_w = int(px_w * scale)
    pic_h = int(px_h * scale)
    left = int((prs.slide_width - pic_w) / 2)
    top = int((prs.slide_height - pic_h) / 2)
    slide.shapes.add_picture(str(artifact_path), left, top, width=pic_w, height=pic_h)


def build_base_pptx(rows, out_pptx):
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)
    blank = prs.slide_layouts[6]
    for row in rows:
        slide = prs.slides.add_slide(blank)
        add_full_slide_picture(slide, row["path"], prs)
    out_pptx.parent.mkdir(parents=True, exist_ok=True)
    prs.save(out_pptx)


def next_rid(xml):
    nums = [int(num) for num in re.findall(r'Id="rId(\d+)"', xml)]
    return f"rId{max(nums, default=0) + 1}"


def add_relationship(xml, rel_type, target):
    if rel_type in xml and target in xml:
        return xml
    rel = f'<Relationship Id="{next_rid(xml)}" Type="{rel_type}" Target="{target}"/>'
    return xml.replace("</Relationships>", rel + "</Relationships>")


def add_override(xml, part_name, content_type):
    if f'PartName="{part_name}"' in xml:
        return xml
    override = f'<Override PartName="{part_name}" ContentType="{content_type}"/>'
    return xml.replace("</Types>", override + "</Types>")


def notes_slide_xml(note_text):
    paras = []
    for line in note_text.split("\n"):
        if not line:
            paras.append("<a:p/>")
        else:
            paras.append(
                '<a:p><a:r><a:rPr lang="en-US" sz="1200"/>'
                f"<a:t>{escape(line)}</a:t></a:r></a:p>"
            )
    return f'''<?xml version='1.0' encoding='UTF-8' standalone='yes'?>
<p:notes xmlns:a="http://schemas.openxmlformats.org/drawingml/2006/main" xmlns:p="http://schemas.openxmlformats.org/presentationml/2006/main" xmlns:r="http://schemas.openxmlformats.org/officeDocument/2006/relationships"><p:cSld><p:spTree><p:nvGrpSpPr><p:cNvPr id="1" name=""/><p:cNvGrpSpPr/><p:nvPr/></p:nvGrpSpPr><p:grpSpPr><a:xfrm><a:off x="0" y="0"/><a:ext cx="0" cy="0"/><a:chOff x="0" y="0"/><a:chExt cx="0" cy="0"/></a:xfrm></p:grpSpPr><p:sp><p:nvSpPr><p:cNvPr id="2" name="Slide Image Placeholder 1"/><p:cNvSpPr><a:spLocks noGrp="1"/></p:cNvSpPr><p:nvPr><p:ph type="sldImg" idx="2"/></p:nvPr></p:nvSpPr><p:spPr/></p:sp><p:sp><p:nvSpPr><p:cNvPr id="3" name="Notes Placeholder 2"/><p:cNvSpPr><a:spLocks noGrp="1"/></p:cNvSpPr><p:nvPr><p:ph type="body" idx="3" sz="quarter"/></p:nvPr></p:nvSpPr><p:spPr/><p:txBody><a:bodyPr/><a:lstStyle/>{"".join(paras)}</p:txBody></p:sp><p:sp><p:nvSpPr><p:cNvPr id="4" name="Slide Number Placeholder 3"/><p:cNvSpPr><a:spLocks noGrp="1"/></p:cNvSpPr><p:nvPr><p:ph type="sldNum" idx="5" sz="quarter"/></p:nvPr></p:nvSpPr><p:spPr/></p:sp></p:spTree></p:cSld><p:clrMapOvr><a:masterClrMapping/></p:clrMapOvr></p:notes>'''


def notes_rels_xml(slide_idx):
    return f'''<?xml version='1.0' encoding='UTF-8' standalone='yes'?>
<Relationships xmlns="{REL_NS}"><Relationship Id="rId1" Type="{NOTES_MASTER_REL}" Target="../notesMasters/notesMaster1.xml"/><Relationship Id="rId2" Type="{SLIDE_REL}" Target="../slides/slide{slide_idx}.xml"/></Relationships>'''


def load_notes_master(template_pptx):
    if not template_pptx or not template_pptx.exists():
        return None
    try:
        with ZipFile(template_pptx) as zf:
            master = zf.read("ppt/notesMasters/notesMaster1.xml")
            rels = zf.read("ppt/notesMasters/_rels/notesMaster1.xml.rels")
            theme = zf.read("ppt/theme/theme2.xml") if "ppt/theme/theme2.xml" in zf.namelist() else None
            return master, rels, theme
    except KeyError:
        return None


def add_notes_to_pptx(pptx_path, rows, template_pptx=None, notes_parts=None):
    if notes_parts is None:
        notes_parts = load_notes_master(template_pptx)
    if notes_parts is None:
        raise RuntimeError("A template PPTX with notes master is required to add speaker notes.")
    notes_master_xml, notes_master_rels_xml, notes_theme_xml = notes_parts

    with ZipFile(pptx_path, "r") as zin:
        package = {name: zin.read(name) for name in zin.namelist()}

    package["ppt/notesMasters/notesMaster1.xml"] = notes_master_xml
    package["ppt/notesMasters/_rels/notesMaster1.xml.rels"] = notes_master_rels_xml
    if notes_theme_xml is not None:
        package["ppt/theme/theme2.xml"] = notes_theme_xml

    content_types = package["[Content_Types].xml"].decode("utf-8")
    content_types = add_override(content_types, "/ppt/notesMasters/notesMaster1.xml", NOTES_MASTER_CONTENT_TYPE)
    if notes_theme_xml is not None:
        content_types = add_override(content_types, "/ppt/theme/theme2.xml", THEME_CONTENT_TYPE)
    for idx in range(1, len(rows) + 1):
        content_types = add_override(content_types, f"/ppt/notesSlides/notesSlide{idx}.xml", NOTES_SLIDE_CONTENT_TYPE)
    package["[Content_Types].xml"] = content_types.encode("utf-8")

    pres_rels_name = "ppt/_rels/presentation.xml.rels"
    pres_rels = package[pres_rels_name].decode("utf-8")
    pres_rels = add_relationship(pres_rels, NOTES_MASTER_REL, "notesMasters/notesMaster1.xml")
    package[pres_rels_name] = pres_rels.encode("utf-8")

    for idx, row in enumerate(rows, 1):
        note_text = f"Pathname: {row['path'].resolve()}\nCaption:\n{row['caption']}"
        package[f"ppt/notesSlides/notesSlide{idx}.xml"] = notes_slide_xml(note_text).encode("utf-8")
        package[f"ppt/notesSlides/_rels/notesSlide{idx}.xml.rels"] = notes_rels_xml(idx).encode("utf-8")
        slide_rels_name = f"ppt/slides/_rels/slide{idx}.xml.rels"
        slide_rels = package[slide_rels_name].decode("utf-8")
        slide_rels = add_relationship(slide_rels, NOTES_SLIDE_REL, f"../notesSlides/notesSlide{idx}.xml")
        package[slide_rels_name] = slide_rels.encode("utf-8")

    with ZipFile(pptx_path, "w", ZIP_DEFLATED) as zout:
        for name, data in package.items():
            zout.writestr(name, data)


def write_checklist(rows, out_csv):
    out_csv.parent.mkdir(parents=True, exist_ok=True)
    with out_csv.open("w", newline="") as handle:
        writer = csv.DictWriter(
            handle,
            fieldnames=[
                "slug",
                "kind",
                "status",
                "scientific_claim_clear",
                "axes_legends_readable",
                "caption_matches_result",
                "needed_edits",
                "artifact_path",
                "svg_path",
                "caption",
            ],
        )
        writer.writeheader()
        for row in rows:
            writer.writerow(
                {
                    "slug": row["slug"],
                    "kind": row["kind"],
                    "status": "draft",
                    "scientific_claim_clear": "",
                    "axes_legends_readable": "",
                    "caption_matches_result": "",
                    "needed_edits": "",
                    "artifact_path": row["path"].relative_to(ROOT),
                    "svg_path": row["svg_path"].relative_to(ROOT) if row["svg_path"] else "",
                    "caption": row["caption"],
                }
            )


def main():
    parser = argparse.ArgumentParser(description="Build the figure/table review deck")
    parser.add_argument("--out-dir", default=ROOT / "paper-writing/figure_review", type=Path)
    parser.add_argument("--template-pptx", default=None, type=Path)
    args = parser.parse_args()

    out_dir = args.out_dir
    out_pptx = out_dir / "figure_review_deck.pptx"
    out_csv = out_dir / "figure_review_checklist.csv"
    template = args.template_pptx or out_pptx

    rows = collect_artifacts(out_dir)
    notes_parts = load_notes_master(template)
    if notes_parts is None:
        raise RuntimeError(f"Could not load notes master from template deck: {template}")

    with tempfile.TemporaryDirectory() as tmp:
        base_pptx = Path(tmp) / "figure_review_base.pptx"
        build_base_pptx(rows, base_pptx)
        out_dir.mkdir(parents=True, exist_ok=True)
        out_pptx.write_bytes(base_pptx.read_bytes())

    add_notes_to_pptx(out_pptx, rows, notes_parts=notes_parts)
    write_checklist(rows, out_csv)

    print(f"Wrote {out_pptx}")
    print(f"Wrote {out_csv}")
    print(f"Slides: {len(rows)}")
    for idx, row in enumerate(rows, 1):
        print(f"{idx:02d} {row['kind']} {row['slug']} -> {row['path'].name}")


if __name__ == "__main__":
    main()
